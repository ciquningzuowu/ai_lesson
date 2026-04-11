"""辅助工具模块

核心功能：
- 文档解析工具（PDF、PPT）
- 文本处理工具
- 响应格式化工具
- 异步任务管理

优化：
- 使用 aiofiles 实现异步文件操作
- 支持多种文档格式
- 提供统一的错误处理机制
"""

import re
import json
from typing import Optional, Any, Dict, Union
from pathlib import Path
import aiofiles
import aiofiles.os

from models.content_processing import ContentType, ParseContentResponse
from utils.node_monitor import node_state


# ============= 文档解析工具 =============

class DocumentParser:
    """文档解析器

    支持 PDF、PPT、TXT 等格式的文档解析
    """

    @staticmethod
    async def parse_text(file_path: str) -> str:
        """解析纯文本文件

        Args:
            file_path: 文件路径

        Returns:
            文本内容
        """
        node_state("helpers.document", "parse_text_file_01", phase="enter", extra={"path": file_path})
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            content = await f.read()
        node_state("helpers.document", "parse_text_file_02", phase="exit", extra={"chars": len(content)})
        return content

    @staticmethod
    async def parse_pdf(file_path: str) -> str:
        """解析 PDF 文件

        Args:
            file_path: 文件路径

        Returns:
            提取的文本内容
        """
        node_state("helpers.document", "parse_pdf_01", phase="enter", extra={"path": file_path})
        try:
            from pypdf import PdfReader
        except ImportError:
            node_state("helpers.document", "parse_pdf_skip", phase="error", message="未安装 pypdf")
            return ""

        text_parts = []
        reader = PdfReader(file_path)
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
        node_state(
            "helpers.document",
            "parse_pdf_02",
            phase="exit",
            message="PDF 页文本提取完成",
            extra={"pages": len(reader.pages), "chars": sum(len(t) for t in text_parts)},
        )
        return "\n".join(text_parts)

    @staticmethod
    async def extract_pdf_images(file_path: str) -> list[dict]:
        """提取 PDF 中的图片信息

        Args:
            file_path: 文件路径

        Returns:
            图片信息列表，每项包含 image_index, page_num, description
        """
        node_state("helpers.document", "extract_pdf_images_01", phase="enter", extra={"path": file_path})
        try:
            from pypdf import PdfReader
            from pypdf.errors import PdfReadError
        except ImportError:
            node_state("helpers.document", "extract_pdf_images_skip", phase="error", message="未安装 pypdf")
            return []

        images_info = []
        try:
            reader = PdfReader(file_path)
            image_count = 0

            for page_num, page in enumerate(reader.pages, 1):
                if "/XObject" in page["/Resources"]:
                    x_objects = page["/Resources"]["/XObject"].get_object()
                    for obj in x_objects:
                        if x_objects[obj]["/Subtype"] == "/Image":
                            image_count += 1
                            images_info.append({
                                "image_index": image_count,
                                "page_num": page_num,
                                "description": f"PDF第{page_num}页的图片 {image_count}",
                                "width": x_objects[obj].get("/Width", 0),
                                "height": x_objects[obj].get("/Height", 0),
                            })

            node_state(
                "helpers.document",
                "extract_pdf_images_02",
                phase="exit",
                extra={"pages": len(reader.pages), "images": len(images_info)},
            )
        except Exception as e:
            node_state("helpers.document", "extract_pdf_images_error", phase="error", message=str(e))

        return images_info

    @staticmethod
    async def describe_pdf_image(image_path: str, llm_client=None) -> str:
        """使用 LLM 描述 PDF 中的图片内容

        Args:
            image_path: 图片路径
            llm_client: LLM 客户端（可选）

        Returns:
            图片描述
        """
        try:
            import base64
            from utils.llm_client import async_generate

            with open(image_path, "rb") as img_file:
                img_data = base64.b64encode(img_file.read()).decode("utf-8")

            prompt = f"""请描述这张图片的内容。如果这是图表或图示，请详细说明其中的数据、趋势或关键信息。

图片数据（base64）: {img_data[:100]}...

请用简洁的语言描述图片内容："""

            description = await async_generate(prompt=prompt)
            return description
        except Exception as e:
            node_state("helpers.document", "describe_pdf_image_error", phase="error", message=str(e))
            return f"图片描述生成失败: {str(e)}"

    @staticmethod
    async def parse_pdf_with_images(file_path: str, describe_images: bool = True) -> dict:
        """解析 PDF 文件并提取图片信息

        Args:
            file_path: 文件路径
            describe_images: 是否描述图片内容

        Returns:
            包含文本和图片信息的字典
        """
        node_state("helpers.document", "parse_pdf_with_images_01", phase="enter", extra={"path": file_path})

        text_content = await DocumentParser.parse_pdf(file_path)
        images_info = await DocumentParser.extract_pdf_images(file_path)

        result = {
            "text": text_content,
            "images": images_info,
            "image_count": len(images_info),
        }

        if describe_images and images_info:
            descriptions = []
            for img in images_info:
                desc = await DocumentParser._get_image_description(img)
                if desc:
                    descriptions.append(f"[图片{img['image_index']}]: {desc}")
            result["image_descriptions"] = descriptions

        node_state(
            "helpers.document",
            "parse_pdf_with_images_02",
            phase="exit",
            extra={"text_chars": len(text_content), "images": len(images_info)},
        )
        return result

    @staticmethod
    async def _get_image_description(image_info: dict) -> str:
        """获取图片描述

        Args:
            image_info: 图片信息

        Returns:
            图片描述
        """
        return f"第{image_info['page_num']}页图片 (尺寸: {image_info.get('width', '?')}x{image_info.get('height', '?')})"

    @staticmethod
    async def parse_ppt(file_path: str) -> str:
        """解析 PPT 文件

        仅支持 .pptx 格式（Office 2007+）

        Args:
            file_path: 文件路径

        Returns:
            提取的文本内容
        """
        import os
        import logging
        logger = logging.getLogger(__name__)

        node_state("helpers.document", "parse_ppt_01", phase="enter", extra={"path": file_path})

        # 验证文件存在
        if not os.path.exists(file_path):
            logger.error(f"[parse_ppt] 文件不存在: {file_path}")
            raise FileNotFoundError(f"文件不存在: {file_path}")

        # 验证文件格式（必须是 .pptx）
        if not file_path.lower().endswith('.pptx'):
            raise ValueError(
                f"不支持的文件格式: {os.path.splitext(file_path)[1]}。"
                "仅支持 .pptx 格式（Office 2007+）。"
                "请将您的 .ppt 文件另存为 .pptx 格式后再上传。"
            )

        logger.info(f"[parse_ppt] 开始解析 PPTX 文件: {file_path}")

        try:
            from pptx import Presentation
        except ImportError:
            node_state("helpers.document", "parse_ppt_skip", phase="error", message="未安装 python-pptx")
            raise ImportError("请安装 python-pptx 库: pip install python-pptx")

        text_parts = []
        prs = Presentation(file_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

                # 正确检查 shape 是否包含表格（避免 python-pptx 异常）
                try:
                    if hasattr(shape, "table") and shape.table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip() if hasattr(cell, "text") else "")
                            if any(row_text):
                                text_parts.append(" | ".join(row_text))
                except Exception:
                    # 如果不是表格 shape，跳过
                    pass

        joined = "\n".join(text_parts)
        node_state(
            "helpers.document",
            "parse_ppt_02",
            phase="exit",
            extra={"slides": len(prs.slides), "chars": len(joined)},
        )
        return joined

    @classmethod
    async def parse_file(
        cls,
        file_path: str,
        file_type: ContentType,
    ) -> str:
        """根据文件类型解析文档

        Args:
            file_path: 文件路径
            file_type: 文件类型

        Returns:
            提取的文本内容
        """
        node_state(
            "helpers.document",
            "parse_file_dispatch",
            phase="enter",
            message="按类型分发解析",
            extra={"file_type": str(file_type), "path": file_path},
        )
        parsers = {
            ContentType.TEXT: cls.parse_text,
            ContentType.PDF: cls.parse_pdf,
            ContentType.PPT: cls.parse_ppt,
        }

        parser = parsers.get(file_type)
        if not parser:
            node_state("helpers.document", "parse_file_unsupported", phase="error", extra={"file_type": str(file_type)})
            return ""
        out = await parser(file_path)
        node_state(
            "helpers.document",
            "parse_file_done",
            phase="exit",
            extra={"file_type": str(file_type), "out_len": len(out) if out else 0},
        )
        return out

    @classmethod
    async def parse_from_content(
        cls,
        content: Union[str, bytes],
        file_type: ContentType,
    ) -> str:
        """从内容直接解析

        Args:
            content: 文件内容或文本
            file_type: 文件类型

        Returns:
            处理后的文本
        """
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="ignore")

        # 如果是 TEXT 类型，直接返回或做基本清理
        if file_type == ContentType.TEXT:
            return cls.clean_text(content)

        # 其他类型内容尝试基本处理
        return cls.clean_text(content)

    @staticmethod
    def clean_text(text: str) -> str:
        """清理文本

        - 去除多余空白
        - 标准化换行
        - 去除特殊字符

        Args:
            text: 输入文本

        Returns:
            清理后的文本
        """
        # 去除多余空白
        text = re.sub(r"\s+", " ", text)
        # 标准化换行
        text = re.sub(r"[\r\n]+", "\n", text)
        # 去除首尾空白
        text = text.strip()
        return text


# ============= 文本处理工具 =============

class TextProcessor:
    """文本处理器

    提供文本处理功能
    """

    @staticmethod
    def truncate_text(text: str, max_length: int = 1000, suffix: str = "...") -> str:
        """截断文本

        Args:
            text: 输入文本
            max_length: 最大长度
            suffix: 截断后缀

        Returns:
            截断后的文本
        """
        if len(text) <= max_length:
            return text
        return text[: max_length - len(suffix)] + suffix

    @staticmethod
    def extract_key_sentences(text: str, num_sentences: int = 5) -> list[str]:
        """提取关键句子

        简单实现：通过句子长度和标点符号判断重要性

        Args:
            text: 输入文本
            num_sentences: 提取数量

        Returns:
            关键句子列表
        """
        # 按句号、问号、感叹号分割
        sentences = re.split(r"[。！？\.\!\?]", text)
        sentences = [s.strip() for s in sentences if s.strip()]

        # 按长度排序（通常较长的句子信息量更大）
        sentences.sort(key=len, reverse=True)

        return sentences[:num_sentences]

    @staticmethod
    def normalize_whitespace(text: str) -> str:
        """规范化空白字符

        Args:
            text: 输入文本

        Returns:
            处理后的文本
        """
        return re.sub(r"\s+", " ", text).strip()

    @staticmethod
    def remove_special_chars(text: str, keep_punctuation: bool = True) -> str:
        """移除特殊字符

        Args:
            text: 输入文本
            keep_punctuation: 是否保留标点

        Returns:
            处理后的文本
        """
        if keep_punctuation:
            return re.sub(r"[^\w\s\u4e00-\u9fff。！？，、：；""''（）《》【】\.\!\?\,\;\:\"\'\(\)\<\>]+", "", text)
        return re.sub(r"[^\w\s\u4e00-\u9fff]+", "", text)


# ============= 响应格式化工具 =============

class ResponseFormatter:
    """响应格式化工具"""

    @staticmethod
    def success_response(data: Any, message: str = "Success") -> dict:
        """成功响应格式

        Args:
            data: 响应数据
            message: 消息

        Returns:
            格式化后的响应
        """
        return {
            "status": "success",
            "message": message,
            "data": data,
        }

    @staticmethod
    def error_response(message: str, code: str = "ERROR", details: Any = None) -> dict:
        """错误响应格式

        Args:
            message: 错误消息
            code: 错误代码
            details: 详细错误信息

        Returns:
            ��式化后的响应
        """
        response = {
            "status": "error",
            "message": message,
            "code": code,
        }
        if details:
            response["details"] = details
        return response

    @staticmethod
    def paginated_response(
        items: list,
        page: int = 1,
        page_size: int = 20,
        total: int = 0,
    ) -> dict:
        """分页响应格式

        Args:
            items: 数据项列表
            page: 当前页码
            page_size: 每页数量
            total: 总数

        Returns:
            格式化后的响应
        """
        return {
            "status": "success",
            "data": {
                "items": items,
                "pagination": {
                    "page": page,
                    "page_size": page_size,
                    "total": total,
                    "total_pages": (total + page_size - 1) // page_size if page_size > 0 else 0,
                },
            },
        }


# ============= 异步任务管理 =============

class TaskManager:
    """异步任务管理器

    管理后台任务的执行和状态追踪
    """

    def __init__(self):
        self._tasks: Dict[str, Any] = {}
        self._results: Dict[str, Any] = {}

    def register_task(self, task_id: str, task: Any):
        """注册任务

        Args:
            task_id: 任务ID
            task: 异步任务对象
        """
        self._tasks[task_id] = task

    def get_task(self, task_id: str) -> Optional[Any]:
        """获取任务"""
        return self._tasks.get(task_id)

    def set_result(self, task_id: str, result: Any):
        """设置任务结果"""
        self._results[task_id] = result

    def get_result(self, task_id: str) -> Optional[Any]:
        """获取任务结果"""
        return self._results.get(task_id)

    def remove_task(self, task_id: str):
        """移除任务"""
        self._tasks.pop(task_id, None)
        self._results.pop(task_id, None)


# ============= JSON 工具 =============

def safe_json_loads(text: str, default: Any = None) -> Any:
    """安全的 JSON 解析

    Args:
        text: JSON 文本
        default: 默认值

    Returns:
        解析后的对象或默认值
    """
    try:
        return json.loads(text)
    except (json.JSONDecodeError, TypeError):
        return default


def safe_json_dumps(obj: Any, default: str = "{}") -> str:
    """安全的 JSON 序列化

    Args:
        obj: 对象
        default: 默认值

    Returns:
        JSON 字符串或默认值
    """
    try:
        return json.dumps(obj, ensure_ascii=False)
    except (TypeError, ValueError):
        return default


# ============= 文件工具 =============

async def ensure_dir(path: str):
    """确保目录存在"""
    Path(path).mkdir(parents=True, exist_ok=True)


async def file_exists(path: str) -> bool:
    """检查文件是否存在"""
    return await aiofiles.os.path.exists(path)


async def read_file(path: str, encoding: str = "utf-8") -> str:
    """异步读取文件"""
    async with aiofiles.open(path, "r", encoding=encoding) as f:
        return await f.read()


async def write_file(path: str, content: str, encoding: str = "utf-8"):
    """异步写入文件"""
    await ensure_dir(str(Path(path).parent))
    async with aiofiles.open(path, "w", encoding=encoding) as f:
        await f.write(content)


# 全局任务管理器
_task_manager = TaskManager()


def get_task_manager() -> TaskManager:
    """获取全局任务管理器"""
    return _task_manager


__all__ = [
    "DocumentParser",
    "TextProcessor",
    "ResponseFormatter",
    "TaskManager",
    "safe_json_loads",
    "safe_json_dumps",
    "ensure_dir",
    "file_exists",
    "read_file",
    "write_file",
    "get_task_manager",
]
